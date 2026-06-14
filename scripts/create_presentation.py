# -*- coding: utf-8 -*-
"""
HAAT_DA_Presentation.pptx  —  Refined Executive Dark
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

OUT = r'C:\Users\Yosef\PycharmProjects\HAAT_assignment\output\HAAT_DA_Presentation.pptx'

# ─────────────────────────────────────────────
# פלטת צבעים — Refined Executive Dark
# ─────────────────────────────────────────────
BG          = RGBColor(0x11, 0x18, 0x27)   # gray-900 — רקע ראשי
CARD        = RGBColor(0x1F, 0x29, 0x37)   # gray-800 — כרטיסים
CARD2       = RGBColor(0x27, 0x35, 0x4A)   # blue-gray — כרטיסים כהים
BORDER      = RGBColor(0x37, 0x41, 0x51)   # gray-700 — קווי גבול

AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # amber-500 — צבע עיקרי
AMBER_L     = RGBColor(0xFC, 0xD3, 0x4D)   # amber-300 — הדגשות
BLUE        = RGBColor(0x60, 0xA5, 0xFA)   # blue-400
GREEN       = RGBColor(0x34, 0xD3, 0x99)   # emerald-400
RED         = RGBColor(0xF8, 0x71, 0x71)   # red-400
ORANGE      = RGBColor(0xFB, 0x92, 0x3C)   # orange-400

WHITE       = RGBColor(0xF9, 0xFA, 0xFB)   # near-white
GREY        = RGBColor(0x9C, 0xA3, 0xAF)   # gray-400 — טקסט משני
GREY_D      = RGBColor(0x37, 0x41, 0x51)   # gray-700 — הפרדות

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = float(prs.slide_width)
H = float(prs.slide_height)

# ─────────────────────────────────────────────
# כלי עיצוב
# ─────────────────────────────────────────────

def S():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    sl.shapes._spTree.remove(bg._element)
    sl.shapes._spTree.insert(2, bg._element)
    return sl

def R(sl, l, t, w, h, fill, line_c=None, lw=1.5, rot=0):
    shp = sl.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_c:
        shp.line.color.rgb = line_c
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    if rot: shp.rotation = rot
    return shp

def O(sl, l, t, w, h, fill, line_c=None, lw=1.0):
    shp = sl.shapes.add_shape(9,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_c:
        shp.line.color.rgb = line_c
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def T(sl, text, l, t, w, h,
      sz=14, bold=False, color=WHITE,
      align=PP_ALIGN.RIGHT, italic=False,
      rtl=True, font='Arial'):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        pp = p._p.get_or_add_pPr()
        pp.set('rtl', '1')
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = font
    return tb

def ML(sl, items, l, t, w, h,
       d_sz=13, d_color=WHITE,
       align=PP_ALIGN.RIGHT, sp=4, font='Arial'):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if   isinstance(item, str):          tx,sz,bd,cl = item,d_sz,False,d_color
        elif len(item)==1:                   tx,sz,bd,cl = item[0],d_sz,False,d_color
        elif len(item)==2:                   tx,sz,bd,cl = item[0],item[1],False,d_color
        elif len(item)==3:                   tx,sz,bd,cl = item[0],item[1],item[2],d_color
        else:                                tx,sz,bd,cl = item
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.alignment = align
        if True:
            pp = p._p.get_or_add_pPr()
            pp.set('rtl', '1')
        p.space_before = Pt(sp)
        r = p.add_run()
        r.text = tx; r.font.size = Pt(sz); r.font.bold = bd
        r.font.color.rgb = cl; r.font.name = font
    return tb

# ─────────────────────────────────────────────
# קומפוננטות UI
# ─────────────────────────────────────────────

HEADER_H   = 1.42   # גובה פס כותרת
FOOTER_Y   = 6.95   # תחתית

def header(sl, tag, title, num, total=7):
    """פס כותרת עם background ברור, tag + כותרת גדולה + מספר שקף"""
    # רקע ה-header
    R(sl, 0, 0, 13.33, HEADER_H, CARD)
    # קו amber עבה בתחתית ה-header
    R(sl, 0, HEADER_H - 0.065, 13.33, 0.065, AMBER)

    # תג (קטן, amber) — ימין
    T(sl, tag, 0.5, 0.10, 12.33, 0.30,
      sz=10, color=AMBER, align=PP_ALIGN.RIGHT)
    # כותרת ראשית — גדולה, לבנה, ברורה
    T(sl, title, 0.5, 0.42, 11.8, 0.82,
      sz=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # מספר שקף — עיגול amber שמאל
    O(sl, 0.28, 0.38, 0.82, 0.82, AMBER)
    T(sl, str(num), 0.28, 0.38, 0.82, 0.82,
      sz=24, bold=True, color=BG, align=PP_ALIGN.CENTER, rtl=False)

    # footer
    R(sl, 0, FOOTER_Y, 13.33, 0.04, AMBER)
    R(sl, 0, FOOTER_Y + 0.04, 13.33, 0.51, CARD)
    T(sl, 'HAAT Delivery  ·  ניתוח נתונים 2024',
      0.4, FOOTER_Y + 0.06, 11.0, 0.38,
      sz=9, color=GREY, align=PP_ALIGN.RIGHT)
    T(sl, f'{num} / {total}', 12.0, FOOTER_Y + 0.06, 1.1, 0.38,
      sz=10, bold=True, color=AMBER,
      align=PP_ALIGN.LEFT, rtl=False)

def card(sl, l, t, w, h, accent=None, bg=CARD):
    """כרטיס עם רקע ובורדר שמאלי אופציונלי"""
    R(sl, l, t, w, h, bg, line_c=BORDER, lw=0.5)
    if accent:
        R(sl, l, t, 0.07, h, accent)

def num_card(sl, l, t, w, h, number, label, sub='', nc=AMBER, accent=None):
    """כרטיס KPI עם מספר גדול"""
    card(sl, l, t, w, h, accent=accent or nc, bg=CARD)
    T(sl, number, l+0.12, t+0.15, w-0.24, h*0.50,
      sz=44, bold=True, color=nc,
      align=PP_ALIGN.CENTER, rtl=False)
    R(sl, l+0.35, t+h*0.61, w-0.7, 0.03, BORDER)
    T(sl, label, l+0.12, t+h*0.65, w-0.24, h*0.24,
      sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        T(sl, sub, l+0.12, t+h*0.85, w-0.24, h*0.16,
          sz=10, color=GREY, align=PP_ALIGN.CENTER)

def info_card(sl, l, t, w, h, title, items,
              tc=AMBER, accent=AMBER, d_sz=12.5, bg=CARD):
    """כרטיס תוכן עם כותרת ופסקאות"""
    card(sl, l, t, w, h, accent=accent, bg=bg)
    cy = t + 0.18
    if title:
        T(sl, title, l+0.25, cy, w-0.38, 0.42,
          sz=13.5, bold=True, color=tc, align=PP_ALIGN.RIGHT)
        R(sl, l+0.25, cy+0.44, w-0.38, 0.025, GREY_D)
        cy += 0.58
    ML(sl, items, l+0.25, cy, w-0.38, h-(cy-t)-0.18,
       d_sz=d_sz, align=PP_ALIGN.RIGHT, sp=5)

def step(sl, l, t, num, text, sub='', c=AMBER):
    """שלב ממוספר"""
    card(sl, l, t, 7.85, 1.10, accent=c, bg=CARD)
    O(sl, l+0.18, t+0.27, 0.56, 0.56, c)
    T(sl, str(num), l+0.18, t+0.27, 0.56, 0.56,
      sz=16, bold=True, color=BG,
      align=PP_ALIGN.CENTER, rtl=False)
    T(sl, text, l+0.96, t+0.10, 6.65, 0.48,
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    if sub:
        T(sl, sub, l+0.96, t+0.58, 6.65, 0.42,
          sz=12, color=GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# שקופית 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
s1 = S()

# אלמנטים גיאומטריים דקורטיביים
R(s1, -4.0, 4.2, 10.0, 10.0, CARD2, rot=32)   # כרטיס ענק מסובב ימין תחתון
R(s1, -3.2, 4.9,  8.0,  8.0, CARD,  rot=32)

O(s1,  9.8, -1.5, 6.5, 6.5, CARD2)             # עיגול גדול פינה שמאל-עליון
O(s1, 10.4, -0.9, 5.3, 5.3, BG)                # חתוך פנים

O(s1,  0.3,  5.8, 0.9, 0.9, AMBER)             # נקודות amber קטנות
O(s1,  1.55, 6.0, 0.5, 0.5, AMBER)
O(s1,  2.35, 6.15, 0.3, 0.3, AMBER)

# תוכן מרכזי
T(s1, 'H A A T   D E L I V E R Y', 2.3, 1.60, 8.73, 0.50,
  sz=12, color=AMBER, bold=False,
  align=PP_ALIGN.CENTER, rtl=False, font='Arial')

R(s1, 2.3, 2.12, 8.73, 0.055, AMBER)

T(s1, 'ניתוח נתוני ביצועים', 2.3, 2.22, 8.73, 1.05,
  sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s1, 'פברואר – מאי 2024', 2.3, 3.25, 8.73, 0.60,
  sz=26, bold=False, color=AMBER_L, align=PP_ALIGN.CENTER)

R(s1, 2.3, 3.90, 8.73, 0.055, BORDER)

T(s1, '28,645 הזמנות   ·   11 אזורים   ·   918 שליחים   ·   4 חודשים',
  2.3, 4.04, 8.73, 0.48, sz=12, color=GREY, align=PP_ALIGN.CENTER)
T(s1, 'ממצא מרכזי  ·  3 מספרים  ·  ממצא מפתיע  ·  לא הוסבר  ·  המלצה  ·  שאלה',
  2.3, 4.60, 8.73, 0.40, sz=11, color=GREY_D, align=PP_ALIGN.CENTER)

# footer title slide
R(s1, 0, FOOTER_Y, 13.33, 0.04, AMBER)
R(s1, 0, FOOTER_Y + 0.04, 13.33, 0.51, CARD)
T(s1, 'HAAT Delivery  ·  ניתוח נתונים 2024',
  0.4, FOOTER_Y + 0.06, 11.0, 0.38, sz=9, color=GREY, align=PP_ALIGN.RIGHT)
T(s1, '1 / 7', 12.0, FOOTER_Y + 0.06, 1.1, 0.38,
  sz=10, bold=True, color=AMBER, align=PP_ALIGN.LEFT, rtl=False)


# ═══════════════════════════════════════════════════════════════════
# שקופית 2 — הממצא החשוב ביותר
# ═══════════════════════════════════════════════════════════════════
s2 = S()
header(s2, 'שקופית   1 / 6  —  הממצא החשוב ביותר',
       '1 מתוך 3 הזמנות לא מגיעה ללקוח', 2)

CONTENT_Y = HEADER_H + 0.20
CONTENT_H = FOOTER_Y - CONTENT_Y - 0.10

# ציטוט hero
R(s2, 0.5, CONTENT_Y, 8.0, 0.055, AMBER)
T(s2, '❝  זו לא בעיה של ביקוש — זו בעיה של היצע שליחים  ❞',
  0.5, CONTENT_Y + 0.07, 8.0, 0.62,
  sz=16.5, bold=True, color=AMBER_L, align=PP_ALIGN.RIGHT, italic=True)

# 3 כרטיסי KPI
num_card(s2, 0.50, CONTENT_Y + 0.82, 3.70, 1.75,
         '66.1%', 'שיעור מסירה כולל', '18,929 הזמנות נמסרו',
         nc=GREEN, accent=GREEN)
num_card(s2, 4.45, CONTENT_Y + 0.82, 3.70, 1.75,
         '9,716', 'הזמנות שלא נמסרו', 'הסתיימו ללא מסירה',
         nc=RED, accent=RED)
num_card(s2, 0.50, CONTENT_Y + 2.75, 7.65, 1.75,
         '₪1,159,000', 'פוטנציאל הכנסה שלא מומש', 'ב-18 שבועות בלבד  ·  AOV ₪119.44',
         nc=ORANGE, accent=ORANGE)

# כרטיס ניתוח — ימין
info_card(s2, 8.48, CONTENT_Y + 0.82, 4.45, 3.68,
          title='ראיה: שבת vs ראשון',
          accent=BLUE,
          items=[
    ('שבת — יום שיא ביקוש', 13, True, WHITE),
    ('271 הזמנות ביום, מסירה 63.1% בלבד ↓', 12, False, GREY),
    ('', 3, False, WHITE),
    ('ראשון — ביקוש נמוך ב-19%', 13, True, WHITE),
    ('227 הזמנות ביום, מסירה 67.5% ↑', 12, False, GREY),
    ('', 3, False, WHITE),
    ('כל 100 הזמנות נוספות', 12.5, True, AMBER_L),
    ('= 4.4 מסירות שנאבדות', 12.5, True, AMBER_L),
    ('', 3, False, WHITE),
    ('כשאין שליחים — ההזמנות נופלות', 12, True, GREEN),
])


# ═══════════════════════════════════════════════════════════════════
# שקופית 3 — 3 מספרים
# ═══════════════════════════════════════════════════════════════════
s3 = S()
header(s3, 'שקופית   2 / 6  —  שלושה מספרים שמספרים את הסיפור',
       'הנתונים מאחורי הממצא המרכזי', 3)

CY3 = HEADER_H + 0.22
CH3 = FOOTER_Y - CY3 - 0.08

def big3(sl, l, num, nc, tag, label, subs, verdict, vc):
    W3 = 4.02
    card(sl, l, CY3, W3, CH3, accent=nc, bg=CARD)
    # tag
    T(sl, tag, l+0.22, CY3+0.14, W3-0.35, 0.40,
      sz=26, bold=True, color=nc, align=PP_ALIGN.RIGHT, rtl=False)
    # big number
    T(sl, num, l+0.10, CY3+0.52, W3-0.20, 1.30,
      sz=58, bold=True, color=nc, align=PP_ALIGN.CENTER, rtl=False)
    R(sl, l+0.30, CY3+1.85, W3-0.60, 0.03, BORDER)
    T(sl, label, l+0.15, CY3+1.95, W3-0.30, 0.50,
      sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, l+0.30, CY3+2.50, W3-0.60, 0.03, BORDER)
    ML(sl, subs, l+0.22, CY3+2.60, W3-0.35, CH3-3.20,
       d_sz=11.5, d_color=GREY, align=PP_ALIGN.RIGHT, sp=5)
    # verdict
    R(sl, l, CY3+CH3-0.60, W3, 0.60, nc)
    T(sl, verdict, l+0.15, CY3+CH3-0.56, W3-0.25, 0.50,
      sz=11.5, bold=True, color=BG, align=PP_ALIGN.CENTER)

big3(s3, 0.30,  '66.1%', GREEN, '01',
     'שיעור מסירה כולל',
     ['9,716 הזמנות לא הגיעו',
      'AOV ₪119.44',
      '₪1.16M פוטנציאל שאבד'],
     '📌  הבסיס לכל שאר', BG)

big3(s3, 4.65, '63.1%', RED, '02',
     'מסירה בשבת — יום השיא',
     ['vs 67.5% בראשון',
      '19% פחות הזמנות = מסירה טובה',
      'כל 100 הזמנות → 4.4 מסירות פחות'],
     '⚠  שיא ביקוש = שפל מסירה', WHITE)

big3(s3, 9.01, '48.4%', ORANGE, '03',
     'מסירה 8–14 אפריל',
     ['שפל מוחלט בכל 18 השבועות',
      'הנפח דווקא עלה ב-15.4%',
      'חופף לתקיפת איראן 13.4.24'],
     '🔴  מחייב חקירה', BG)


# ═══════════════════════════════════════════════════════════════════
# שקופית 4 — ממצא מפתיע
# ═══════════════════════════════════════════════════════════════════
s4 = S()
header(s4, 'שקופית   3 / 6  —  ממצא שהפתיע אותי',
       'לקוחות המזומן — לא מי שחשבנו', 4)

CY4 = HEADER_H + 0.22

# מספר ענק + context — חצי שמאל
T(s4, '57.6%', 0.40, CY4+0.10, 6.10, 1.80,
  sz=92, bold=True, color=AMBER,
  align=PP_ALIGN.CENTER, rtl=False)

R(s4, 0.70, CY4+1.94, 5.50, 0.04, AMBER)
T(s4, 'מלקוחות המזומן', 0.40, CY4+2.04, 6.10, 0.50,
  sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s4, 'רשמו כרטיס אשראי — ובחרו לשלם מזומן בכל זאת', 0.40, CY4+2.52, 6.10, 0.45,
  sz=14, color=GREY, align=PP_ALIGN.CENTER)

# badge
R(s4, 1.10, CY4+3.15, 4.30, 0.70, CARD2, line_c=AMBER, lw=1.0)
T(s4, '= 9,190 לקוחות מוכנים להמרה',
  1.15, CY4+3.22, 4.20, 0.55,
  sz=14, bold=True, color=AMBER_L, align=PP_ALIGN.CENTER)

# ניתוח ימין
info_card(s4, 6.88, CY4, 6.05, FOOTER_Y-CY4-0.08,
          title='מה זה אומר?', accent=BLUE, tc=AMBER_L,
          items=[
    ('הם כבר סומכים על HAAT ✓', 13.5, True, WHITE),
    ('רשמו כרטיס = מחויבות גבוהה', 12.5, False, GREY),
    ('', 4, False, WHITE),
    ('הם בוחרים מזומן בגלל:', 13.5, True, WHITE),
    ('→  הרגל / פרטיות / חוסר תמריץ', 12.5, False, GREY),
    ('', 4, False, WHITE),
    ('מה ניתן לעשות?', 13.5, True, AMBER_L),
    ('→  קמפיין קאשבק 5% על תשלום ראשון', 12.5, False, WHITE),
    ('→  הנחה בהזמנה הבאה בכרטיס', 12.5, False, WHITE),
    ('→  תכנית נאמנות מבוססת אשראי', 12.5, False, WHITE),
    ('', 4, False, WHITE),
    ('הסגמנט הקל ביותר להמרה —', 12.5, True, GREEN),
    ('הוא כבר קיים, רק צריך תמריץ', 12.5, True, GREEN),
])


# ═══════════════════════════════════════════════════════════════════
# שקופית 5 — מה שלא הוסבר
# ═══════════════════════════════════════════════════════════════════
s5 = S()
header(s5, 'שקופית   4 / 6  —  מה שלא הצלחתי להסביר',
       'נפח עלה 15% — אבל המסירה קרסה ל-48.4%', 5)

CY5 = HEADER_H + 0.22

# כותרת משנה
T(s5, 'שבוע 8–14 אפריל 2024:', 0.50, CY5, 7.0, 0.45,
  sz=15, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

# טבלה
COL_L = [0.50, 3.68, 5.62, 7.45, 9.10]
COL_W = [3.10, 1.85, 1.75, 1.55, 1.65]
COL_H = ['שבוע', 'הזמנות', 'שינוי', 'מסירה', 'מצב']

# header row
R(s5, 0.50, CY5+0.52, 10.38, 0.52, CARD2)
for txt, ll, ww in zip(COL_H, COL_L, COL_W):
    T(s5, txt, ll+0.08, CY5+0.56, ww-0.12, 0.40,
      sz=12, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)

rows = [
    ('01–07 אפריל', '1,363', '—',         '62.1%', '🟡 תקין', False, CARD),
    ('08–14 אפריל', '1,573', '+15.4% ↑',  '48.4% ↓','🔴 משבר', True,  RGBColor(0x2A, 0x10, 0x10)),
    ('15–21 אפריל', '1,987', '+26.3% ↑',  '67.4%', '🟢 חזרה', False, CARD),
]
for ri, (d, o, ch, dr, st, is_crisis, bg) in enumerate(rows):
    y = CY5 + 1.10 + ri * 0.65
    R(s5, 0.50, y, 10.38, 0.60, bg)
    if is_crisis:
        R(s5, 0.50, y, 0.07, 0.60, RED)
    vals = [d, o, ch, dr, st]
    crs  = [WHITE if not is_crisis else RED,
            GREY  if not is_crisis else WHITE,
            ORANGE if (is_crisis) else GREY,
            RED   if is_crisis else (GREEN if ri==2 else GREY),
            RED   if is_crisis else (GREEN if ri==2 else GREY)]
    for vv, ll, ww, cc in zip(vals, COL_L, COL_W, crs):
        T(s5, vv, ll+0.08, y+0.11, ww-0.12, 0.40,
          sz=12.5, bold=is_crisis, color=cc, align=PP_ALIGN.RIGHT)

# אזהרה
R(s5, 0.50, CY5+3.15, 10.38, 0.70, CARD2, line_c=ORANGE, lw=1.0)
T(s5, '13 אפריל 2024 — תקיפת טילים ומל"טים מאיראן',
  0.65, CY5+3.20, 10.08, 0.34,
  sz=13, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
T(s5, 'שליחים בחרו שלא לצאת → עלייה בביקוש + קריסה בהיצע = 48.4%',
  0.65, CY5+3.50, 10.08, 0.30,
  sz=12, color=GREY, align=PP_ALIGN.RIGHT)

# כרטיס "מה חסר"
info_card(s5, 10.95, CY5, 2.00, FOOTER_Y-CY5-0.08,
          title='מה חסר?', accent=ORANGE, tc=ORANGE, d_sz=11,
          bg=CARD2,
          items=[
    ('לוגי פעילות שליחים', 11, False, WHITE),
    ('כמה היו online?', 10.5, False, GREY),
    ('', 3, False, WHITE),
    ('גיאוגרפיה של הכשלים', 11, False, WHITE),
    ('', 3, False, WHITE),
    ('תוכנית חירום', 11, True, ORANGE),
    ('לאירועי ביטחון?', 11, True, ORANGE),
])


# ═══════════════════════════════════════════════════════════════════
# שקופית 6 — המלצה
# ═══════════════════════════════════════════════════════════════════
s6 = S()
header(s6, 'שקופית   5 / 6  —  המלצה שניתן לבצע השבוע הבא',
       'גייסו שליחים ייעודיים לסוף שבוע', 6)

CY6 = HEADER_H + 0.22

step(s6, 0.42, CY6+0.05, 1,
     'גייסו 20–30 שליחים ייעודיים לשישי + שבת',
     'באזורים 1, 4 ו-10 — ימות השיא', GREEN)
step(s6, 0.42, CY6+1.28, 2,
     'הציעו בונוס שיפט — +10% על כל מסירה בסוף שבוע',
     'עלות: ~₪2,000 | תמריץ אמיתי ומדיד', BLUE)
step(s6, 0.42, CY6+2.51, 3,
     'הפעילו כבר השבת הקרובה',
     'יש נתון Baseline ברור — ניתן למדוד מיד', AMBER)
step(s6, 0.42, CY6+3.74, 4,
     'מדדו ודווחו — שיעור מסירה שבת לפני vs אחרי',
     'ה-KPI פשוט, ברור ואין ויכוח על ההצלחה', ORANGE)

# ROI card — ימין
card(s6, 8.60, CY6+0.05, 4.35, 5.28, accent=GREEN, bg=CARD2)
T(s6, 'ROI מחושב 💰', 8.85, CY6+0.22, 3.90, 0.42,
  sz=14, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
R(s6, 8.85, CY6+0.68, 3.90, 0.03, BORDER)
ML(s6, [
    ('+5% מסירה בשבת =', 14, True, WHITE),
    ('~68 הזמנות נוספות', 13, False, GREY),
    ('', 5, False, WHITE),
    ('AOV ₪122.76 →', 14, True, WHITE),
    ('₪8,348 הכנסה נוספת', 13, False, GREY),
    ('', 5, False, WHITE),
    ('עלות בונוסים: ~₪2,000', 13, False, GREY),
], 8.85, CY6+0.80, 3.90, 2.65, align=PP_ALIGN.RIGHT, sp=5)

R(s6, 8.60, CY6+3.60, 4.35, 1.73, GREEN)
T(s6, 'ROI × 4', 8.72, CY6+3.72, 4.10, 0.78,
  sz=38, bold=True, color=BG,
  align=PP_ALIGN.CENTER, rtl=False)
T(s6, 'מהשבת הראשונה ✓', 8.72, CY6+4.46, 4.10, 0.40,
  sz=13, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# שקופית 7 — שאלה שלא נשאלה
# ═══════════════════════════════════════════════════════════════════
s7 = S()
header(s7, 'שקופית   6 / 6  —  שאלה שלא נשאלה — אבל צריך לשאול',
       'מהו ה-LTV של לקוח חוזר לעומת לקוח חד-פעמי?', 7)

CY7 = HEADER_H + 0.22

# מספר ענק + label
T(s7, '42.3%', 0.40, CY7+0.10, 6.00, 1.70,
  sz=88, bold=True, color=BLUE,
  align=PP_ALIGN.CENTER, rtl=False)

R(s7, 0.70, CY7+1.85, 5.40, 0.04, BLUE)
T(s7, 'שיעור לקוחות חוזרים', 0.40, CY7+1.96, 6.00, 0.50,
  sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s7, 'מסך כלל הלקוחות בדאטאסט', 0.40, CY7+2.44, 6.00, 0.40,
  sz=13, color=GREY, align=PP_ALIGN.CENTER)

info_card(s7, 0.40, CY7+3.00, 6.00, FOOTER_Y-CY7-3.10,
          title='מה שאנחנו לא יודעים:', accent=BLUE, tc=BLUE,
          items=[
    ('→  האם לקוח חוזר שווה יותר כלכלית?', 12.5, False, WHITE),
    ('→  מה גורם ללקוח לחזור?', 12.5, False, WHITE),
    ('→  האם מסירה כושלת = לקוח שלא חוזר?', 12.5, False, WHITE),
    ('→  מהו ה-Churn Rate לפי אזור?', 12.5, False, WHITE),
])

# פאנל ימין
info_card(s7, 6.78, CY7, 6.15, FOOTER_Y-CY7-0.08,
          title='למה זה קריטי:', accent=AMBER, tc=AMBER_L,
          bg=CARD2,
          items=[
    ('עלות לקוח חדש גבוהה פי 5–7', 13.5, True, WHITE),
    ('מעלות שימור לקוח קיים', 13, False, GREY),
    ('', 5, False, WHITE),
    ('אם מסירה כושלת = לקוח לא חוזר:', 13.5, True, RED),
    ('ה-ROI האמיתי של שיפור delivery rate', 13, False, WHITE),
    ('גבוה בהרבה ממה שחישבנו', 13, False, GREY),
    ('', 5, False, WHITE),
    ('הנתון הנדרש:', 13.5, True, AMBER_L),
    ('Repeat Purchase Rate פר קוהורט', 13, False, WHITE),
    ('לפי אזור  ·  עסק  ·  שיטת תשלום', 12.5, False, GREY),
])

# closing banner
R(s7, 0, FOOTER_Y-0.70, 13.33, 0.70, CARD2)
R(s7, 0, FOOTER_Y-0.70, 13.33, 0.04, AMBER)
T(s7, '❝  HAAT יושבת על ביקוש אמיתי. האתגר: לא לאבד 1 מכל 3 לקוחות בדרך. השליח הוא הנקודה — והוא גם ההזדמנות  ❞',
  0.40, FOOTER_Y-0.66, 12.53, 0.62,
  sz=12.5, bold=True, color=AMBER_L, align=PP_ALIGN.RIGHT, italic=True)

# ═══════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f'Saved: {OUT}  ({len(prs.slides)} slides)')
